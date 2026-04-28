from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

ROOT = Path(r"G:\SYSTEM_DESIGN\RuoYi-Vue")
OUTPUT_PPT = ROOT / "毕业设计中期检查汇报_符合真实进度.pptx"

def set_font(run, size, bold=False, color=(34, 34, 34), name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Cm(1.2), Cm(0.8), Cm(22), Cm(1.8))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, 24, bold=True, color=(24, 24, 24))

    if subtitle:
        sub_box = slide.shapes.add_textbox(Cm(1.25), Cm(2.0), Cm(22), Cm(0.8))
        sub_tf = sub_box.text_frame
        p = sub_tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        set_font(r, 10.5, color=(110, 110, 110))

def add_bg(slide, accent=(46, 125, 50)):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(0.45), Cm(19.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*accent)
    bar.line.fill.background()

    head = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(25.4), Cm(0.28)
    )
    head.fill.solid()
    head.fill.fore_color.rgb = RGBColor(*accent)
    head.line.fill.background()

def add_bullets(slide, items, left=1.5, top=3.2, width=22, height=10.5, font_size=17):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.level = 0
        p.space_after = Pt(12)
        p.line_spacing = 1.3
        for run in p.runs:
            set_font(run, font_size)
    return box

def add_card(slide, left, top, width, height, title, lines, fill=(247, 250, 247), title_color=(46, 125, 50)):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(222, 230, 222)

    title_box = slide.shapes.add_textbox(Cm(left + 0.4), Cm(top + 0.25), Cm(width - 0.8), Cm(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, 13.5, bold=True, color=title_color)

    content_box = slide.shapes.add_textbox(Cm(left + 0.45), Cm(top + 1.2), Cm(width - 0.9), Cm(height - 1.4))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    for i, line in enumerate(lines):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = "▪ " + line
        p.space_after = Pt(8)
        p.line_spacing = 1.2
        for run in p.runs:
            set_font(run, 11.5, color=(55, 55, 55))

def add_footer(slide, text="毕业论文中期检查汇报"):
    box = slide.shapes.add_textbox(Cm(19.2), Cm(18.0), Cm(5.2), Cm(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    set_font(r, 9.5, color=(120, 120, 120))

def build_presentation():
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]

    # 1 cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    cover = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.3), Cm(2.0), Cm(22.4), Cm(13.5)
    )
    cover.fill.solid()
    cover.fill.fore_color.rgb = RGBColor(248, 251, 247)
    cover.line.color.rgb = RGBColor(225, 233, 223)
    
    title_box = slide.shapes.add_textbox(Cm(2.0), Cm(5.0), Cm(18.8), Cm(4.8))
    tf = title_box.text_frame
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "毕业论文中期检查汇报"
    set_font(r1, 28, bold=True, color=(35, 80, 40))
    p2 = tf.add_paragraph()
    p2.text = "课题：《基于 SpringBoot + Vue 的农产品销售系统》"
    p2.space_before = Pt(14)
    for run in p2.runs:
        set_font(run, 16, color=(80, 80, 80), bold=True)

    info = slide.shapes.add_textbox(Cm(2.0), Cm(11.0), Cm(10.5), Cm(2.2))
    itf = info.text_frame
    for i, line in enumerate(["汇报人：__________", "指导教师：__________", "汇报日期：__________"]):
        p = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
        p.text = line
        p.space_after = Pt(8)
        for run in p.runs:
            set_font(run, 14)

    # 2 已完成的工作 - 系统方面
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "1. 已完成的工作 —— 系统功能开发")
    
    add_card(slide, 1.6, 3.0, 7.0, 11.5, "基础框架与权限", [
        "基于若依前后端分离框架进行二次开发。",
        "明确管理员、卖家、买家三类用户角色。",
        "依托现有体系，实现不同角色的功能权限与数据范围隔离。"
    ])
    add_card(slide, 9.0, 3.0, 7.0, 11.5, "商品与商城展示", [
        "商品管理：实现商品发布、信息维护、库存及上下架管理。",
        "商城展示：完成商品列表前台展示和详情查看。"
    ])
    add_card(slide, 16.4, 3.0, 7.0, 11.5, "交易与统计闭环", [
        "购物车：加入、改量、删除及结算前信息维护。",
        "订单：创建、支付模拟、发收货、取消订单及库存回补。",
        "统计分析：订单、销售和商品数据的基础统计展示。"
    ])
    add_footer(slide)

    # 3 已完成的工作 - 论文方面
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "2. 已完成的工作 —— 论文撰写")
    
    add_card(slide, 1.6, 3.2, 10.6, 10.0, "前期准备与初稿结构", [
        "完成开题报告撰写，题目确定为《基于 SpringBoot + Vue 的农产品销售系统》。",
        "初稿主体内容成型：包含中英文摘要、引言、绪论、相关理论与关键技术。",
        "覆盖完整的工程流程：包含需求分析、总体设计、详细设计、系统测试及总结与展望。"
    ])
    add_card(slide, 12.8, 3.2, 10.6, 10.0, "核心系统实现说明", [
        "结合系统源码，对核心模块的实现逻辑进行了详细的文字说明。",
        "重点说明：商品模块底层设计逻辑。",
        "重点说明：购物车状态流转与存储机制。",
        "重点说明：订单模块的状态机流转处理。",
        "重点说明：权限控制与统计模块的数据聚合逻辑。"
    ])
    add_footer(slide)

    # 4 未完成的工作 - 系统方面
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "3. 未完成的工作 —— 系统完善与测试")
    
    add_bullets(slide, [
        "测试支撑材料补充：需补充基础功能测试用例、测试截图和运行结果记录，确保各角色下主要业务流程的正确性。",
        "多角色协同验证：需进一步检查和强化管理员、卖家、买家不同角色在特定模块（如订单流转、权限隔离）的边界处理。",
        "交互与体验细节优化：重点完善页面交互细节、各类异常情况提示友好度，以及统计数据的可视化展示效果。",
        "核心功能稳定性增强：继续打磨商品浏览、购物车操作、订单完整流转过程，尤其是取消订单时的库存并发回补逻辑。"
    ], left=1.6, top=3.5, width=22.0, height=11.0, font_size=16)
    add_footer(slide)

    # 5 未完成的工作 - 论文方面
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "4. 未完成的工作 —— 论文修订与规范")
    
    add_bullets(slide, [
        "图表绘制与补充：当前缺失核心系统图表，需绘制总体架构图、功能模块图和数据库 E-R 图，以直观展现系统设计。",
        "实证材料补充：需要结合系统的最终测试，补充具体的系统功能截图、真实测试结果截图以及主要的测试用例表。",
        "格式规范与排版：语言表达需更偏向学术化；章节衔接、图表编号、参考文献格式需严格对照本科毕业论文撰写规范进行调整。",
        "终稿定型：根据指导教师的中期检查意见，继续多轮修改论文内容，直至完善形成最终版本。"
    ], left=1.6, top=3.5, width=22.0, height=11.0, font_size=16)
    add_footer(slide)
    
    # 6 总结与致谢
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    
    summary_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(2.0), Cm(4.0), Cm(21.4), Cm(7.0)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(248, 251, 247)
    summary_box.line.color.rgb = RGBColor(225, 233, 223)

    text = slide.shapes.add_textbox(Cm(2.6), Cm(4.8), Cm(20.0), Cm(5.0))
    tf = text.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "阶段总结与展望"
    p1.space_after = Pt(12)
    set_font(p1.runs[0], 20, bold=True, color=(35, 80, 40))
    
    p2 = tf.add_paragraph()
    p2.text = "总体而言，项目目前处于【主体开发完成、进入细节完善与论文润色】阶段，进度符合开题时制定的计划。后续将严格遵照规范，扎实完成测试截图补充、图表绘制与格式排版，确保系统稳定运行和论文高质量定稿。"
    p2.line_spacing = 1.3
    set_font(p2.runs[0], 16, color=(50, 50, 50))
    
    thanks = slide.shapes.add_textbox(Cm(2.0), Cm(14.0), Cm(21.4), Cm(2.0))
    ptf = thanks.text_frame
    p3 = ptf.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r = p3.add_run()
    r.text = "感谢各位老师的聆听与指导！"
    set_font(r, 24, bold=True, color=(46, 125, 50))

    add_footer(slide)

    prs.save(OUTPUT_PPT)

if __name__ == "__main__":
    build_presentation()

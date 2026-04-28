from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

ROOT = Path(r"G:\SYSTEM_DESIGN\RuoYi-Vue")
OUTPUT_PPT = ROOT / "毕业设计中期检查汇报_12页详尽版.pptx"

def set_font(run, size, bold=False, color=(34, 34, 34), name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Cm(1.2), Cm(0.8), Cm(22), Cm(1.8))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, 24, bold=True, color=(24, 24, 24))

    if subtitle:
        sub_box = slide.shapes.add_textbox(Cm(1.25), Cm(2.0), Cm(22), Cm(0.8))
        sub_tf = sub_box.text_frame
        p = sub_tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        set_font(r, 10.5, color=(110, 110, 110))

def add_bg(slide, accent=(46, 125, 50)):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(0.45), Cm(19.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*accent)
    bar.line.fill.background()

    head = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(25.4), Cm(0.28)
    )
    head.fill.solid()
    head.fill.fore_color.rgb = RGBColor(*accent)
    head.line.fill.background()

def add_bullets(slide, items, left=1.5, top=3.2, width=22, height=10.5, font_size=17, line_spacing=1.3):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.level = 0
        p.space_after = Pt(12)
        p.line_spacing = line_spacing
        for run in p.runs:
            set_font(run, font_size)
    return box

def add_card(slide, left, top, width, height, title, lines, fill=(247, 250, 247), title_color=(46, 125, 50)):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(222, 230, 222)

    title_box = slide.shapes.add_textbox(Cm(left + 0.4), Cm(top + 0.25), Cm(width - 0.8), Cm(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, 13.5, bold=True, color=title_color)

    content_box = slide.shapes.add_textbox(Cm(left + 0.45), Cm(top + 1.2), Cm(width - 0.9), Cm(height - 1.4))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    for i, line in enumerate(lines):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = "▪ " + line
        p.space_after = Pt(8)
        p.line_spacing = 1.2
        for run in p.runs:
            set_font(run, 11.5, color=(55, 55, 55))

def add_footer(slide, text="毕业论文中期检查汇报"):
    box = slide.shapes.add_textbox(Cm(19.2), Cm(18.0), Cm(5.2), Cm(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    set_font(r, 9.5, color=(120, 120, 120))

def build_presentation():
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]

    # 第1页：封面
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    cover = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.3), Cm(2.0), Cm(22.4), Cm(13.5))
    cover.fill.solid()
    cover.fill.fore_color.rgb = RGBColor(248, 251, 247)
    cover.line.color.rgb = RGBColor(225, 233, 223)
    
    title_box = slide.shapes.add_textbox(Cm(2.0), Cm(4.0), Cm(18.8), Cm(4.8))
    tf = title_box.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "毕业论文中期检查汇报"
    set_font(p1.runs[0], 28, bold=True, color=(35, 80, 40))
    p2 = tf.add_paragraph()
    p2.text = "基于 SpringBoot + Vue 的农产品销售系统"
    p2.space_before = Pt(14)
    set_font(p2.runs[0], 20, color=(80, 80, 80), bold=True)

    info = slide.shapes.add_textbox(Cm(2.0), Cm(10.0), Cm(15.0), Cm(5.0))
    itf = info.text_frame
    for i, line in enumerate(["学      院：__________", "专      业：__________", "学生姓名：__________", "指导教师：__________"]):
        p = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
        p.text = line
        p.space_after = Pt(10)
        set_font(p.runs[0], 16)

    # 第2页：汇报内容
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "汇报内容大纲")
    add_bullets(slide, [
        "1. 选题背景与研究目标",
        "2. 当前完成情况总览",
        "3. 系统核心功能与架构实现",
        "4. 论文写作进度",
        "5. 后续工作安排"
    ], left=2.0, top=4.5, font_size=20, line_spacing=1.8)
    add_footer(slide)

    # 第3页：选题背景
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "1. 选题背景")
    add_bullets(slide, [
        "销售渠道有限：传统农产品销售依赖线下批发，农户直销渠道不畅。",
        "经营管理粗放：商品库存、订单状态和资金流水缺乏信息化管理工具。",
        "信息透明度不足：买卖双方信息不对称，商品溯源和售后沟通成本高。",
        "痛点解决：迫切需要一个垂直领域的在线销售平台，打破信息壁垒，规范交易流程。"
    ], top=4.0, font_size=18, line_spacing=1.5)
    add_footer(slide)

    # 第4页：研究目标
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "2. 研究目标")
    add_card(slide, 1.6, 3.5, 22.0, 3.5, "技术选型", ["基于 SpringBoot + Vue 以及若依前后端分离框架，构建高性能的 B2C/C2C 农产品销售平台。"])
    add_card(slide, 1.6, 7.5, 22.0, 3.5, "用户角色", ["面向 管理员、卖家、买家 三类角色，提供差异化的操作入口与数据权限。"])
    add_card(slide, 1.6, 11.5, 22.0, 3.5, "核心目标", ["完成商品管理、购物车、订单流转、严格的权限控制及数据统计分析等核心业务模块开发。"])
    add_footer(slide)

    # 第5页：当前完成情况总览
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "3. 当前完成情况总览")
    add_card(slide, 1.6, 4.0, 10.6, 10.0, "系统工程进度：主体已完成", [
        "系统主体核心功能已全部开发完毕并联调通过。",
        "已经形成完整的商品上架、浏览、下单到发货收货的交易闭环。",
        "后端基础鉴权与前后端接口交互稳定正常。"
    ])
    add_card(slide, 12.8, 4.0, 10.6, 10.0, "论文工程进度：初稿已完成", [
        "论文主体结构与内容已撰写完毕（第1章-第6章）。",
        "已结合实际源码详细说明了系统的核心实现过程。",
        "后续仅需重点补充部分图表、系统测试截图及格式调整。"
    ], fill=(236, 248, 237), title_color=(35, 100, 40))
    add_footer(slide)

    # 第6页：系统功能模块
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "4. 系统已完成功能模块")
    add_card(slide, 1.5, 3.5, 7.0, 5.0, "商品管理", ["商品发布", "商品信息维护", "库存管理", "上下架管理"])
    add_card(slide, 9.0, 3.5, 7.0, 5.0, "商城展示", ["前台商品列表", "商品详情查看"])
    add_card(slide, 16.5, 3.5, 7.0, 5.0, "购物车管理", ["加入购物车", "数量修改", "购物项删除", "结算前信息维护"])
    
    add_card(slide, 1.5, 9.5, 7.0, 5.0, "订单管理", ["订单创建与支付模拟", "卖家发货，买家收货", "订单取消", "库存安全回补"])
    add_card(slide, 9.0, 9.5, 7.0, 5.0, "权限控制", ["基于若依体系构建", "多角色访问隔离", "角色数据范围限制"])
    add_card(slide, 16.5, 9.5, 7.0, 5.0, "统计分析", ["订单数据统计", "销售额数据分析", "商品热销展示"])
    add_footer(slide)

    # 第7页：系统总体架构
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "5. 系统总体架构设计")
    add_bullets(slide, [
        "前端表示层：采用 Vue2 + Element UI，实现响应式组件与友好的后台管理/商城前台界面。",
        "后端逻辑层：基于 Spring Boot 框架，提供高并发下的 RESTful API 接口支持。",
        "安全与权限层：依托 Spring Security 实现严格的登录认证及接口级别的权限拦截。",
        "数据持久层：使用 MyBatis 框架实现高效的 ORM 映射及复杂业务 SQL 处理。",
        "数据存储层：业务数据持久化至 MySQL 数据库，同时引入 Redis 缓存提升频繁访问数据的读取效率。"
    ], top=4.0, font_size=17, line_spacing=1.5)
    add_footer(slide)

    # 第8页：核心业务流程
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "6. 核心业务流程 (交易闭环)")
    add_card(slide, 1.6, 3.5, 22.0, 11.0, "完整订单流转过程", [
        "第一步：买家浏览商品，选择规格后加入购物车。",
        "第二步：买家在购物车内修改数量、勾选商品，点击结算进入订单确认。",
        "第三步：提交订单并触发模拟支付，系统锁定库存。",
        "第四步：卖家后台查看到已支付订单，进行发货操作，订单状态变更为“待收货”。",
        "第五步：买家前台确认收货，交易完成。",
        "异常分支：买家主动取消订单，系统触发库存回补机制，保证库存数据一致性。"
    ], fill=(250, 250, 250), title_color=(40, 40, 40))
    add_footer(slide)

    # 第9页：核心实现重点
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "7. 核心技术实现重点")
    add_card(slide, 1.6, 3.5, 7.0, 11.0, "订单状态机流转", [
        "设计了严谨的订单状态枚举（待支付、待发货、待收货、已完成、已取消）。",
        "后端通过状态校验防止非法流转（如已发货订单不可被买家单方面取消）。"
    ])
    add_card(slide, 9.0, 3.5, 7.0, 11.0, "角色数据范围控制", [
        "基于 AOP 切面与数据权限注解，实现了底层 SQL 的动态拼接。",
        "保证卖家只能看到自己的商品和订单，买家只能操作自己的购物车，管理员拥有全局监管权。"
    ])
    add_card(slide, 16.4, 3.5, 7.0, 11.0, "库存扣减与回补机制", [
        "采用乐观锁/悲观锁机制防止高并发下的库存超卖。",
        "订单取消时，采用事务保证订单状态变更与库存回补的原子性操作。"
    ])
    add_footer(slide)

    # 第10页：基础运行验证
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "8. 基础运行与测试验证")
    add_bullets(slide, [
        "验证方式：目前已通过注册管理员、卖家、买家测试账号，进行全流程的手工模拟操作验证。",
        "功能覆盖：已验证商品维护、商城浏览、购物车加购、订单流转、取消订单触发库存回补、数据统计查看等核心链路。",
        "验证结论：各模块交互正常，数据落地无误，权限隔离生效，主体业务逻辑走通。",
        "后续完善：当前以基础手工验证为主，下一步将补充系统化的测试用例表格及详细的页面运行截图放入论文。"
    ], top=4.0, font_size=18, line_spacing=1.6)
    add_footer(slide)

    # 第11页：论文写作进度
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "9. 论文写作进度详情")
    add_card(slide, 1.6, 3.5, 22.0, 11.0, "主体章节撰写状态", [
        "【已完成】摘要与引言：明确了课题背景与意义。",
        "【已完成】理论与技术：总结了 SpringBoot、Vue、MySQL 等底层框架的作用。",
        "【已完成】需求与设计：完成了系统用例分析与六大模块的详细设计撰写。",
        "【已完成】系统实现：结合源码，着重剖析了购物车、订单和权限模块的具体代码实现。",
        "【已完成】系统测试：基于手工验证情况，初步撰写了测试章节骨架。",
        "【已完成】参考文献：已整理 11 篇高质量中文参考文献。"
    ])
    add_footer(slide)

    # 第12页：后续工作安排
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "10. 后续工作安排")
    add_bullets(slide, [
        "图表绘制：重点补充系统总体架构图、功能模块图、数据库 E-R 图，提升论文直观性。",
        "截图与用例：截取关键业务流程的系统运行图，并整理测试用例表补充至测试章节。",
        "系统细节优化：继续修复在手工测试中发现的页面交互瑕疵，优化异常情况的弹窗提示。",
        "格式规范审查：严格按照学校本科毕业论文撰写规范，调整目录、字体、图表编号及参考文献格式。",
        "确保目标：按时、保质、保量地完成论文终稿，为正式答辩做好准备。"
    ], top=4.0, font_size=18, line_spacing=1.6)
    add_footer(slide)

    prs.save(OUTPUT_PPT)

if __name__ == "__main__":
    build_presentation()

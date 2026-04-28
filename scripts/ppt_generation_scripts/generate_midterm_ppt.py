from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(r"G:\SYSTEM_DESIGN\RuoYi-Vue")
OUTPUT_PPT = ROOT / "农产品销售平台_中期检查汇报.pptx"
OUTPUT_NOTES = ROOT / "农产品销售平台_中期检查讲稿.txt"


def set_font(run, size, bold=False, color=(34, 34, 34), name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Cm(1.2), Cm(0.8), Cm(22), Cm(1.8))
    title_tf = title_box.text_frame
    p = title_tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, 24, bold=True, color=(24, 24, 24))

    if subtitle:
        sub_box = slide.shapes.add_textbox(Cm(1.25), Cm(2.0), Cm(22), Cm(0.8))
        sub_tf = sub_box.text_frame
        p = sub_tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        set_font(r, 10.5, color=(110, 110, 110))


def add_bg(slide, accent=(46, 125, 50)):
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(0.45), Cm(19.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*accent)
    bar.line.fill.background()

    head = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(25.4), Cm(0.28)
    )
    head.fill.solid()
    head.fill.fore_color.rgb = RGBColor(*accent)
    head.line.fill.background()


def add_bullets(slide, items, left=1.5, top=3.2, width=10.6, height=10.5, font_size=17):
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.line_spacing = 1.25
        for run in p.runs:
            set_font(run, font_size)
    return box


def add_card(slide, left, top, width, height, title, lines, fill=(247, 250, 247), title_color=(46, 125, 50)):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(222, 230, 222)

    title_box = slide.shapes.add_textbox(Cm(left + 0.4), Cm(top + 0.25), Cm(width - 0.8), Cm(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_font(r, 13.5, bold=True, color=title_color)

    content_box = slide.shapes.add_textbox(Cm(left + 0.45), Cm(top + 1.0), Cm(width - 0.9), Cm(height - 1.2))
    ctf = content_box.text_frame
    ctf.word_wrap = True
    for i, line in enumerate(lines):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = line
        p.space_after = Pt(6)
        p.line_spacing = 1.2
        for run in p.runs:
            set_font(run, 11.5, color=(55, 55, 55))


def add_footer(slide, text="中期检查汇报"):
    box = slide.shapes.add_textbox(Cm(19.2), Cm(18.0), Cm(5.2), Cm(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    set_font(r, 9.5, color=(120, 120, 120))


def build_presentation():
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)

    blank = prs.slide_layouts[6]

    # 1 cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    cover = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.3), Cm(2.0), Cm(22.4), Cm(13.5)
    )
    cover.fill.solid()
    cover.fill.fore_color.rgb = RGBColor(248, 251, 247)
    cover.line.color.rgb = RGBColor(225, 233, 223)
    add_title(slide, "基于 SpringBoot + Vue 的农产品销售平台", "毕业设计中期检查汇报")

    box = slide.shapes.add_textbox(Cm(2.0), Cm(5.0), Cm(18.8), Cm(4.8))
    tf = box.text_frame
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "中期汇报内容"
    set_font(r1, 26, bold=True, color=(35, 80, 40))
    p2 = tf.add_paragraph()
    p2.text = "项目背景、系统设计、当前进展、存在问题、下一步计划"
    p2.space_before = Pt(10)
    for run in p2.runs:
        set_font(run, 16, color=(80, 80, 80))

    info = slide.shapes.add_textbox(Cm(2.0), Cm(11.4), Cm(10.5), Cm(2.2))
    itf = info.text_frame
    for i, line in enumerate(["学生姓名：__________", "专业班级：__________", "指导教师：__________"]):
        p = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
        p.text = line
        p.space_after = Pt(8)
        for run in p.runs:
            set_font(run, 14)
    add_footer(slide)

    # 2 background
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "1. 课题背景与研究目标")
    add_bullets(
        slide,
        [
            "传统农产品销售环节多、信息不对称，农户和消费者两端都存在实际痛点。",
            "本课题希望做一个面向农产品销售场景的管理平台，覆盖商品、购物车、订单和数据统计等核心流程。",
            "项目采用前后端分离方式开发，前端使用 Vue，后端使用 Spring Boot，便于模块划分和后续维护。",
            "中期阶段的重点不是把所有功能做满，而是先把主要业务链路跑通。"
        ],
        left=1.6,
        top=3.2,
        width=14.3,
        height=10.8,
        font_size=17,
    )
    add_card(
        slide, 17.0, 4.0, 6.0, 4.4, "本课题想解决的问题",
        [
            "商品展示和管理不方便",
            "订单处理流程不完整",
            "不同角色操作边界不清晰"
        ]
    )
    add_card(
        slide, 17.0, 9.0, 6.0, 4.4, "目前明确的目标",
        [
            "完成基础交易闭环",
            "实现角色权限控制",
            "提供基础数据看板"
        ]
    )
    add_footer(slide)

    # 3 design
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "2. 系统总体设计")
    add_card(slide, 1.6, 3.2, 6.8, 4.8, "前端层", ["Vue 2 + Element UI", "负责页面展示、表单交互、接口调用"])
    add_card(slide, 9.2, 3.2, 6.8, 4.8, "后端层", ["Spring Boot + Spring Security", "负责接口、权限、业务逻辑"])
    add_card(slide, 16.8, 3.2, 6.0, 4.8, "数据层", ["MyBatis + MySQL + Redis", "负责数据持久化与缓存"])
    add_card(
        slide, 1.6, 9.0, 21.2, 5.6, "当前系统的主要业务模块",
        [
            "1）商品管理：商品信息维护、上下架、库存管理",
            "2）购物车管理：买家添加商品、修改数量、提交订单准备",
            "3）订单管理：支付、发货、收货、取消等状态流转",
            "4）权限与统计：管理员、卖家、买家三类角色，以及首页经营数据展示"
        ],
        fill=(250, 250, 250),
        title_color=(40, 40, 40),
    )
    add_footer(slide)

    # 4 completed
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "3. 中期阶段已完成的工作")
    add_card(slide, 1.6, 3.0, 6.8, 5.1, "后端部分", [
        "完成商品、购物车、订单等核心接口开发",
        "实现订单创建、库存扣减、状态修改等业务逻辑",
        "完成基于角色的数据访问控制"
    ])
    add_card(slide, 9.2, 3.0, 6.8, 5.1, "前端部分", [
        "完成登录、首页、商品、购物车、订单等页面",
        "实现前后端接口联调",
        "完成基础表格、表单和弹窗交互"
    ])
    add_card(slide, 16.8, 3.0, 6.0, 5.1, "数据库与环境", [
        "整理项目数据库脚本",
        "配置 MySQL、Redis、Maven、Node 环境",
        "完成本地运行和基础测试"
    ])
    add_card(slide, 1.6, 9.2, 21.2, 5.6, "阶段性结果", [
        "目前系统主流程已经能够跑通：商品浏览 -> 加入购物车 -> 提交订单 -> 支付/发货/收货。",
        "管理员、卖家、买家三类用户在系统中的可见菜单和可操作数据存在区分。",
        "首页已经具备基础统计功能，可以展示销售额、订单数量、热销商品等信息。"
    ], fill=(247, 250, 247))
    add_footer(slide)

    # 5 progress
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "4. 当前进度说明")
    add_card(slide, 1.6, 3.3, 4.8, 9.5, "需求分析", ["已完成", "已明确系统角色、功能范围和主要业务流程"], fill=(236, 248, 237))
    add_card(slide, 7.0, 3.3, 4.8, 9.5, "系统设计", ["已完成", "完成模块划分、数据库设计和接口结构整理"], fill=(236, 248, 237))
    add_card(slide, 12.4, 3.3, 4.8, 9.5, "核心功能开发", ["基本完成", "主流程可运行，细节体验和异常处理仍在完善"], fill=(255, 248, 230), title_color=(160, 110, 20))
    add_card(slide, 17.8, 3.3, 5.0, 9.5, "论文与材料", ["进行中", "已形成开题报告，正在整理中期检查材料和论文初稿"], fill=(255, 248, 230), title_color=(160, 110, 20))
    add_footer(slide)

    # 6 problems
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "5. 目前存在的问题")
    add_bullets(
        slide,
        [
            "对 Spring Boot、Vue 这类框架的理解还在加深，很多功能最开始是边学边做。",
            "项目虽然主流程已经通了，但个别细节还不够完善，比如分类统计、页面细节优化、异常提示统一等。",
            "支付和物流部分目前还是偏模拟，离真实商业系统还有差距。",
            "论文部分需要尽快把系统实现内容转成规范文字，避免后期时间过于集中。"
        ],
        left=1.6,
        top=3.4,
        width=14.2,
        height=10.6,
        font_size=17,
    )
    add_card(slide, 17.1, 4.0, 5.8, 4.2, "自我判断", [
        "进度总体可控",
        "系统主体已经成型",
        "后续重点是完善和总结"
    ], fill=(250, 246, 238), title_color=(130, 90, 20))
    add_card(slide, 17.1, 9.0, 5.8, 4.2, "接下来要补的点", [
        "论文内容整理",
        "细节测试与修正",
        "答辩材料准备"
    ], fill=(250, 246, 238), title_color=(130, 90, 20))
    add_footer(slide)

    # 7 plan
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "6. 下一步计划")
    add_card(slide, 1.6, 3.2, 6.8, 10.0, "功能完善", [
        "继续补充和优化商品、订单、统计模块细节",
        "完善页面交互和异常提示",
        "补做必要的测试和运行截图"
    ])
    add_card(slide, 9.2, 3.2, 6.8, 10.0, "论文推进", [
        "按照系统分析、设计、实现、测试的结构整理正文",
        "把已有项目内容转化为论文表达",
        "同步修改格式、图表和参考文献"
    ])
    add_card(slide, 16.8, 3.2, 6.0, 10.0, "答辩准备", [
        "准备项目演示流程",
        "梳理技术点和业务逻辑讲解",
        "提前准备老师可能提问的问题"
    ])
    add_footer(slide)

    # 8 summary
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_title(slide, "7. 阶段总结")
    summary_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.7), Cm(3.5), Cm(21.0), Cm(8.8)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(248, 251, 247)
    summary_box.line.color.rgb = RGBColor(225, 233, 223)

    text = slide.shapes.add_textbox(Cm(2.4), Cm(4.4), Cm(19.0), Cm(6.8))
    tf = text.text_frame
    paras = [
        "目前项目已经完成了中期阶段应有的主要工作，系统核心功能基本具备，能够体现课题的主要设计思路。",
        "接下来我的重点会放在两个方面：一是继续把系统细节完善好，二是把论文和汇报材料整理扎实。",
        "整体来看，本课题能够按计划继续推进，后续目标是顺利完成论文、演示和正式答辩。"
    ]
    for i, line in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(12)
        p.line_spacing = 1.35
        for run in p.runs:
            set_font(run, 19 if i == 0 else 17, bold=(i == 0), color=(48, 48, 48))

    thanks = slide.shapes.add_textbox(Cm(2.4), Cm(13.6), Cm(8), Cm(1))
    ptf = thanks.text_frame
    p = ptf.paragraphs[0]
    r = p.add_run()
    r.text = "谢谢各位老师指导"
    set_font(r, 20, bold=True, color=(35, 80, 40))
    add_footer(slide)

    prs.save(OUTPUT_PPT)


def build_notes():
    notes = """《基于 SpringBoot+Vue 的农产品销售平台》中期检查讲稿

第1页 封面
各位老师好，我的课题是《基于 SpringBoot+Vue 的农产品销售平台设计与实现》。下面我从课题背景、系统设计、目前进展、存在问题以及下一步计划几个方面做一个中期汇报。

第2页 课题背景与研究目标
这个课题主要是针对农产品销售过程中的一些实际问题，比如销售环节多、信息不够透明、管理方式比较分散等。我希望通过这个项目做一个比较完整的农产品销售平台，把商品管理、购物车、订单处理和数据统计这些核心功能串起来。中期阶段我的目标不是把系统做得特别复杂，而是先把主要流程跑通。

第3页 系统总体设计
系统采用前后端分离架构。前端使用 Vue 和 Element UI，负责页面展示和交互；后端使用 Spring Boot，负责接口、权限和业务逻辑；数据层使用 MySQL 和 Redis。当前系统主要包括商品管理、购物车管理、订单管理以及权限与统计四个部分。

第4页 中期阶段已完成的工作
到目前为止，后端已经完成了商品、购物车、订单这些核心接口开发，前端完成了登录、首页、商品、购物车和订单相关页面，并且已经完成前后端联调。数据库脚本和本地开发环境也已经整理好。现在系统的主流程基本可以跑通。

第5页 当前进度说明
从整体进度来看，需求分析和系统设计已经完成，核心功能开发基本完成，论文和材料整理还在进行中。也就是说，项目主体已经做出来了，现在更多是在细化、完善和总结。

第6页 目前存在的问题
目前遇到的问题主要有几方面。第一，我对 Spring Boot 和 Vue 的理解还在继续加深，开发过程中有比较明显的学习过程。第二，系统虽然能运行，但一些细节还需要继续优化。第三，像支付和物流这些部分目前还是模拟实现。第四，论文写作需要尽快跟上开发进度。

第7页 下一步计划
下一步我会继续完善系统细节，补充必要的测试和截图，同时把论文正文按规范整理出来。另外也会提前准备好项目演示流程和答辩中可能会被问到的技术点，避免后期时间太赶。

第8页 阶段总结
整体来说，我认为项目已经完成了中期阶段应该完成的主要内容，系统核心功能已经具备，课题可以继续按计划推进。后续重点是把功能做扎实、把论文写规范，为最终答辩做好准备。
"""
    OUTPUT_NOTES.write_text(notes, encoding="utf-8")


if __name__ == "__main__":
    build_presentation()
    build_notes()
